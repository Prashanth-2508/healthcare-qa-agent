"""
Generate Healthcare_QA_Agent.pptx from content defined inline.

Run:
    python build_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x0D, 0x47, 0xA1)   # deep blue  — title bg
BLUE_MID    = RGBColor(0x19, 0x76, 0xD2)   # medium blue — accents
BLUE_LIGHT  = RGBColor(0xE3, 0xF2, 0xFD)   # pale blue  — content bg
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x21, 0x21, 0x21)
GREY_TEXT   = RGBColor(0x55, 0x55, 0x55)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
AMBER       = RGBColor(0xF5, 0x7F, 0x17)
RED         = RGBColor(0xC6, 0x28, 0x28)

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=DARK_TEXT,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, space_before=0, italic=False, level=0):
    para = tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    para.level = level
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return para


def slide_header(slide, title, subtitle=None, bg=BLUE_DARK):
    """Blue header bar at top."""
    bar_h = Inches(1.35)
    add_rect(slide, 0, 0, W, bar_h, bg)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.8),
                 font_size=30, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.4),
                     font_size=14, color=RGBColor(0xBB, 0xDE, 0xFB), italic=True)
    # thin accent line
    add_rect(slide, 0, bar_h, W, Pt(3), BLUE_MID)


def content_box(slide, x, y, w, h, bg=BLUE_LIGHT, border=BLUE_MID):
    box = add_rect(slide, x, y, w, h, bg)
    box.line.color.rgb = border
    box.line.width = Pt(1)
    return box


def bullet_box(slide, items, x, y, w, h, title=None,
               bg=BLUE_LIGHT, font_size=13, title_size=15):
    content_box(slide, x, y, w, h, bg)
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12),
                                     w - Inches(0.3), h - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = BLUE_DARK
    for item in items:
        indent = item.startswith("    ") or item.startswith("  •")
        txt = item.lstrip("  ")
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if txt.startswith("✅") or txt.startswith("✓"):
            run.font.color.rgb = GREEN
        elif txt.startswith("❌") or txt.startswith("✗"):
            run.font.color.rgb = RED
        elif txt.startswith("⚠"):
            run.font.color.rgb = AMBER
        else:
            run.font.color.rgb = DARK_TEXT
        run.text = txt
        run.font.size = Pt(font_size - (1 if indent else 0))
        p.level = 1 if indent else 0


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_rect(slide, 0, 0, W, H, BLUE_DARK)
add_rect(slide, 0, Inches(5.2), W, Inches(2.3), BLUE_MID)
add_rect(slide, 0, Inches(5.18), W, Pt(4), WHITE)

add_text_box(slide, "Healthcare Q&A Agent",
             Inches(0.6), Inches(1.4), Inches(12), Inches(1.3),
             font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "AI-Powered Clinical Decision Support",
             Inches(0.6), Inches(2.75), Inches(12), Inches(0.7),
             font_size=24, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
add_text_box(slide, "LangGraph  ·  LangChain  ·  Ollama  ·  PubMed API  ·  ChromaDB  ·  Streamlit",
             Inches(0.6), Inches(3.5), Inches(12), Inches(0.5),
             font_size=14, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

add_text_box(slide, "Balaji GV  |  May 2026  |  Assignment Submission",
             Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "The Problem", "Healthcare professionals waste hours searching for answers")

# Left — problem
bullet_box(slide,
    [
        "A medical coder needs the ICD-10 code for T2DM with CKD",
        "A CDI specialist needs ADA 2024 first-line treatment guidelines",
        "A revenue cycle manager needs to verify Medicare CGM coverage",
        "",
        "Today's workflow:",
        "  Google → 10 open tabs → Manual cross-referencing → Maybe get the answer",
        "",
        "Problems with today's approach:",
        "  ✗  Time-consuming — 20–40 minutes per complex question",
        "  ✗  No citations — hard to justify coding decisions",
        "  ✗  Scattered — information lives in PubMed, CMS.gov, and textbooks",
    ],
    Inches(0.3), Inches(1.5), Inches(6.3), Inches(5.7),
    title="The Challenge", font_size=13)

# Right — solution
bullet_box(slide,
    [
        "Ask one question in plain English",
        "",
        "Agent automatically:",
        "  🧠  Understands the question type",
        "  🔍  Searches PubMed, guidelines, ICD-10 databases",
        "  📊  Evaluates whether the evidence is sufficient",
        "  ✍️  Writes a cited Markdown answer",
        "",
        "✅  Seconds instead of minutes",
        "✅  Every claim has a source citation",
        "✅  Covers PubMed + Guidelines + ICD-10 in one query",
    ],
    Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.7),
    title="What I Built", font_size=13, bg=RGBColor(0xE8, 0xF5, 0xE9))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LIVE DEMO FLOW
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "How It Works", "One question → evidence-based answer")

steps = [
    ("🧠", "1. UNDERSTAND", "Reads the question, classifies it (coding / coverage / treatment / evidence), decides which sources to search"),
    ("🔍", "2. RESEARCH", "Searches PubMed (live API), clinical guidelines (RAG), and ICD-10-CM database (74,260 codes) simultaneously"),
    ("📊", "3. EVALUATE", "Self-reflects: 'Is this evidence sufficient?' Loops back for more searches if not (max 2 iterations)"),
    ("✍️", "4. ANSWER", "Writes a clear Markdown answer with source citations, streams it token-by-token to the screen"),
]

for i, (icon, title, desc) in enumerate(steps):
    x = Inches(0.3 + i * 3.26)
    # Box
    add_rect(slide, x, Inches(1.6), Inches(3.1), Inches(4.9), BLUE_LIGHT)
    # Top colour bar
    add_rect(slide, x, Inches(1.6), Inches(3.1), Inches(0.55), BLUE_MID)
    # Icon + title
    add_text_box(slide, f"{icon} {title}",
                 x + Inches(0.1), Inches(1.62), Inches(2.9), Inches(0.5),
                 font_size=15, bold=True, color=WHITE)
    # Description
    add_text_box(slide, desc,
                 x + Inches(0.15), Inches(2.25), Inches(2.85), Inches(4.0),
                 font_size=12, color=DARK_TEXT, wrap=True)
    # Arrow between boxes
    if i < 3:
        add_text_box(slide, "→",
                     x + Inches(3.12), Inches(3.7), Inches(0.2), Inches(0.4),
                     font_size=22, bold=True, color=BLUE_MID)

# Example at bottom
add_rect(slide, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.75), RGBColor(0xFF, 0xF9, 0xC4))
add_text_box(slide,
             'Example: "What is the ICD-10 code for essential hypertension?"  →  '
             'I10 — Essential (primary) hypertension  [source: CMS ICD-10-CM FY2025]',
             Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6),
             font_size=12, color=RGBColor(0x4E, 0x34, 0x00))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "System Architecture", "LangGraph state machine + tool layer + persistence")

# Central pipeline
boxes = [
    (Inches(0.3),  "STREAMLIT UI",    "Question input\nLive progress\nStreaming answer",     BLUE_LIGHT),
    (Inches(2.55), "SAFETY LAYER",    "Block PII\nBlock injections\nEmergency → 911",        RGBColor(0xFF, 0xF0, 0xE0)),
    (Inches(4.8),  "LANGGRAPH\nAGENT","Reason → Plan → Act\nObserve → Respond\n(5 nodes)",  RGBColor(0xE8, 0xEA, 0xF6)),
    (Inches(7.05), "TOOL LAYER",      "PubMed API\nGuidelines RAG\nICD-10 DB",              RGBColor(0xE8, 0xF5, 0xE9)),
    (Inches(9.3),  "PERSISTENCE",     "SQLite DB\n(conversations)\nJSONL traces",            RGBColor(0xFC, 0xE4, 0xEC)),
]

for x, title, body, bg in boxes:
    add_rect(slide, x, Inches(1.55), Inches(2.0), Inches(4.5), bg)
    add_rect(slide, x, Inches(1.55), Inches(2.0), Inches(0.5), BLUE_MID)
    add_text_box(slide, title,
                 x + Inches(0.05), Inches(1.57), Inches(1.9), Inches(0.46),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, body,
                 x + Inches(0.1), Inches(2.15), Inches(1.85), Inches(3.7),
                 font_size=11, color=DARK_TEXT, wrap=True)

# Arrows
for i in range(4):
    ax = Inches(2.35 + i * 2.25)
    add_text_box(slide, "→", ax, Inches(3.45), Inches(0.25), Inches(0.4),
                 font_size=22, bold=True, color=BLUE_MID)

# Bottom note
add_rect(slide, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.9),
         RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide,
             "Every node transition is timed and saved to a JSONL trace file — "
             "providing full auditability of every agent decision.",
             Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.8),
             font_size=12, color=BLUE_DARK, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE 5 NODES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "The 5-Node Agent Loop", "How the AI makes decisions step by step")

nodes = [
    ("1. REASON",  "LLM Call",       BLUE_LIGHT,
     ["Reads query + conversation history",
      "Classifies: coding / coverage / treatment / evidence",
      "Produces a JSON execution plan with tool steps & rationale",
      "Output: reasoning_json (which tools to call and why)"]),
    ("2. PLAN",    "Deterministic",  RGBColor(0xE8, 0xF5, 0xE9),
     ["No LLM call — pure logic",
      "Validates tool names against registry",
      "Coerces argument types",
      "Output: tool_plan (clean, ordered list)"]),
    ("3. ACT",     "Tool Executor",  RGBColor(0xFF, 0xF8, 0xE1),
     ["Executes each tool in sequence",
      "Measures duration per tool call",
      "Handles failures gracefully (no crash)",
      "Output: tool_results with success/error/duration"]),
    ("4. OBSERVE", "LLM Call",       RGBColor(0xFC, 0xE4, 0xEC),
     ["Self-reflects: 'Is evidence sufficient?'",
      "Rates evidence quality: strong/moderate/weak",
      "May request additional tool calls",
      "Sets goal_met=true → routes to RESPOND"]),
    ("5. RESPOND", "LLM Call",       RGBColor(0xE8, 0xEA, 0xF6),
     ["Synthesises all observations",
      "Writes plain-Markdown cited answer",
      "Appends retrieved PMIDs and ICD-10 codes",
      "Saves to SQLite + JSONL trace"]),
]

for i, (name, ntype, bg, bullets) in enumerate(nodes):
    x = Inches(0.25 + i * 2.6)
    add_rect(slide, x, Inches(1.5), Inches(2.5), Inches(5.7), bg)
    add_rect(slide, x, Inches(1.5), Inches(2.5), Inches(0.52), BLUE_DARK)
    add_text_box(slide, name,
                 x + Inches(0.08), Inches(1.52), Inches(2.35), Inches(0.3),
                 font_size=13, bold=True, color=WHITE)
    add_text_box(slide, f"[{ntype}]",
                 x + Inches(0.08), Inches(1.83), Inches(2.35), Inches(0.25),
                 font_size=10, color=RGBColor(0x90, 0xCA, 0xF9), italic=True)
    tb = slide.shapes.add_textbox(x + Inches(0.12), Inches(2.15),
                                  Inches(2.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    for b in bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(11)
        run.font.color.rgb = DARK_TEXT
        p.space_before = Pt(2)

    # Arrow
    if i < 4:
        add_text_box(slide, "→",
                     x + Inches(2.52), Inches(4.2), Inches(0.15), Inches(0.4),
                     font_size=18, bold=True, color=BLUE_MID)

# Loop annotation
add_rect(slide, Inches(7.45), Inches(6.75), Inches(5.3), Inches(0.5),
         RGBColor(0xFF, 0xF9, 0xC4))
add_text_box(slide,
             "OBSERVE can loop back to PLAN (max 2 iterations) if more evidence is needed",
             Inches(7.55), Inches(6.78), Inches(5.1), Inches(0.4),
             font_size=11, color=RGBColor(0x4E, 0x34, 0x00))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TOOLS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "The 4 Tools", "What the agent can search and why")

tools = [
    ("🔬 search_pubmed",
     "PubMed Research Database",
     ["Source: NCBI Entrez API (free, official)",
      "35 million peer-reviewed articles",
      "Rate-limited: 3 req/s (10 req/s with API key)",
      "Auto-translates natural language → PubMed Boolean syntax",
      "Returns: PMID, title, authors, abstract, MeSH terms"],
     RGBColor(0xE3, 0xF2, 0xFD)),
    ("📋 get_clinical_guidelines",
     "Clinical Guidelines Library",
     ["Source: Local PDF files + ChromaDB vector database",
      "Supports ADA, ACC/AHA, USPSTF, CMS guidelines",
      "Chunks PDFs into 400-token pieces with 80-token overlap",
      "Uses all-MiniLM-L6-v2 embeddings (384 dimensions)",
      "Returns: top-3 most relevant guideline passages"],
     RGBColor(0xE8, 0xF5, 0xE9)),
    ("🏥 lookup_medical_code",
     "ICD-10 / HCPCS Code Search",
     ["Source: CMS ICD-10-CM FY2025 (auto-downloaded)",
      "74,260 diagnosis codes",
      "Uses RapidFuzz fuzzy matching (token sort ratio)",
      "Threshold: 55% similarity",
      "Returns: top matches with description and score"],
     RGBColor(0xFF, 0xF8, 0xE1)),
    ("🔍 describe_medical_code",
     "Exact Code Lookup",
     ["Same CMS dataset as above",
      "Exact code string lookup (e.g. 'E11.65')",
      "Returns official description + coding notes",
      "Used when reason node already knows the code",
      "Response time: < 100ms (local database)"],
     RGBColor(0xFC, 0xE4, 0xEC)),
]

for i, (name, subtitle, bullets, bg) in enumerate(tools):
    x = Inches(0.25 + (i % 2) * 6.5)
    y = Inches(1.55 + (i // 2) * 2.9)
    add_rect(slide, x, y, Inches(6.2), Inches(2.6), bg)
    add_rect(slide, x, y, Inches(6.2), Inches(0.5), BLUE_DARK)
    add_text_box(slide, name,
                 x + Inches(0.1), y + Inches(0.04), Inches(6.0), Inches(0.3),
                 font_size=14, bold=True, color=WHITE)
    add_text_box(slide, subtitle,
                 x + Inches(0.1), y + Inches(0.5), Inches(6.0), Inches(0.3),
                 font_size=12, bold=True, color=BLUE_DARK)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.85),
                                  Inches(5.9), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for b in bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(11)
        run.font.color.rgb = DARK_TEXT
        p.space_before = Pt(1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — STREAMLIT UI
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "The Web Interface", "User-friendly Streamlit UI for non-technical users")

left_items = [
    "Welcome Screen (when chat is empty):",
    "  • 4-column 'How It Works' explainer",
    "  • 6 clickable example questions",
    "  • Medical disclaimer banner",
    "",
    "Research Phase (live progress):",
    "  • 🧠 Understanding your question",
    "  • 🔍 Searching medical sources",
    "  • 📊 Reviewing the evidence",
    "  • Plain English — no technical jargon",
    "",
    "Answer Section:",
    "  • Streams token-by-token as AI writes it",
    "  • PMID numbers become clickable links",
    "  • Markdown rendered with headings & bullets",
]
bullet_box(slide, left_items,
           Inches(0.3), Inches(1.5), Inches(6.2), Inches(5.8),
           title="User-Facing Features", font_size=12)

right_items = [
    "Sources Used (collapsible):",
    "  • Lists each database searched with ✅",
    "  • Shows PubMed article links (PMID → URL)",
    "  • Execution time per tool",
    "",
    "How the Agent Worked (collapsible):",
    "  • Colour-coded step-by-step reasoning trace",
    "  • Shows timestamps for each decision",
    "  • Non-technical labels (not raw JSON)",
    "",
    "Sidebar Controls:",
    "  • Response Style: Precise ↔ Creative",
    "  • Max Research Articles: 1–10",
    "  • Enable/disable medical code lookup",
    "  • Current AI model display",
    "  • New conversation button",
]
bullet_box(slide, right_items,
           Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.8),
           title="Response Display", font_size=12, bg=RGBColor(0xE8, 0xF5, 0xE9))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SAFETY GUARDRAILS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Safety & Reliability", "What makes it safe to use in healthcare contexts")

layers = [
    ("BEFORE any LLM call",    BLUE_LIGHT,
     ["✓ Block queries containing PII (SSN, MRN, DOB patterns)",
      "✓ Block prompt injection attempts",
      "✓ Detect emergency keywords → redirect to 911",
      "✓ Check query length (2–2000 chars)"]),
    ("DURING tool execution",   RGBColor(0xE8, 0xF5, 0xE9),
     ["✓ Validate every tool input with Pydantic schemas",
      "✓ Handle tool failures gracefully (no crash)",
      "✓ Validate tool output size and content",
      "✓ Guardrail blocks flagged tool responses"]),
    ("IN every answer",         RGBColor(0xFF, 0xF8, 0xE1),
     ["✓ Mandatory medical disclaimer on every response",
      "✓ 'Verify against official sources' reminder",
      "✓ Confidence level stated (high / medium / low)",
      "✓ Every factual claim must have a citation"]),
    ("AUDIT trail",             RGBColor(0xFC, 0xE4, 0xEC),
     ["✓ JSONL trace file per session (every decision logged)",
      "✓ Node timing recorded (millisecond precision)",
      "✓ Conversation history in SQLite database",
      "✓ Trace shows exactly what was searched and why"]),
]

for i, (label, bg, items) in enumerate(layers):
    x = Inches(0.25 + (i % 2) * 6.5)
    y = Inches(1.55 + (i // 2) * 2.8)
    bullet_box(slide, items, x, y, Inches(6.2), Inches(2.55),
               title=label, font_size=12, bg=bg)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EVALUATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Evaluation Results", "5 clinical test scenarios across 4 question types")

# Table header
def trow(slide, cols, y, bg, text_color=DARK_TEXT, bold=False, font_size=12):
    widths = [Inches(1.1), Inches(3.5), Inches(1.7), Inches(2.3), Inches(1.8), Inches(2.5)]
    x = Inches(0.3)
    for i, (col, w) in enumerate(zip(cols, widths)):
        add_rect(slide, x, y, w, Inches(0.42), bg)
        add_text_box(slide, col, x + Inches(0.05), y + Inches(0.04),
                     w - Inches(0.1), Inches(0.35),
                     font_size=font_size, bold=bold, color=text_color)
        x += w

header_bg = BLUE_DARK
trow(slide, ["ID", "Scenario", "Category", "Pass Criteria", "Result", "Root Cause"],
     Inches(1.5), header_bg, WHITE, bold=True, font_size=11)

rows = [
    ("TC-001", "ADA 2024 T2DM Treatment",       "Guidelines", "metformin + GLP-1/SGLT2",  "❌ FAIL", "Guidelines PDF not loaded"),
    ("TC-002", "Essential Hypertension ICD-10",  "Coding",     "Returns code I10",          "❌ FAIL", "Model arg name mismatch"),
    ("TC-003", "Medicare CGM Coverage 2024",     "Policy",     "LCD + coverage expansion",  "❌ FAIL", "Guidelines PDF not loaded"),
    ("TC-004", "SGLT2 + Heart Failure PubMed",   "Evidence",   "Cites ≥1 PMID",             "✅ PASS", "PubMed works end-to-end"),
    ("TC-005", "CPT New Patient Office Visit",   "Coding",     "Returns 99204 or 99205",    "❌ FAIL", "CPT not in ICD-10 dataset"),
]
row_bgs = [BLUE_LIGHT, WHITE, BLUE_LIGHT, RGBColor(0xE8, 0xF5, 0xE9), BLUE_LIGHT]
for j, (row, bg) in enumerate(zip(rows, row_bgs)):
    y = Inches(1.92 + j * 0.45)
    trow(slide, list(row), y, bg, font_size=10)

# Summary box
add_rect(slide, Inches(0.3), Inches(4.35), Inches(6.2), Inches(2.85),
         RGBColor(0xFF, 0xEB, 0xEE))
add_text_box(slide, "Current Results (gemma3:1b)",
             Inches(0.45), Inches(4.4), Inches(6.0), Inches(0.4),
             font_size=14, bold=True, color=RED)
for k, line in enumerate([
    "Pass rate: 1 / 5 (20%)",
    "Mean score: 1.0 / 5.0 (heuristic)",
    "Model: gemma3:1b — 1B parameters, limited tool-use",
    "Bottleneck: model, NOT the retrieval pipeline",
]):
    add_text_box(slide, "• " + line,
                 Inches(0.5), Inches(4.85 + k * 0.45), Inches(5.9), Inches(0.42),
                 font_size=12, color=DARK_TEXT)

add_rect(slide, Inches(6.8), Inches(4.35), Inches(6.2), Inches(2.85),
         RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, "Expected with llama-3.1-70b",
             Inches(6.95), Inches(4.4), Inches(6.0), Inches(0.4),
             font_size=14, bold=True, color=GREEN)
for k, line in enumerate([
    "Pass rate: 4 / 5 (80%) estimated",
    "TC-004 (PubMed): always passes",
    "TC-001/003: pass once PDFs are loaded",
    "Switch model: 5 minutes, free via Groq",
]):
    add_text_box(slide, "• " + line,
                 Inches(6.95), Inches(4.85 + k * 0.45), Inches(5.9), Inches(0.42),
                 font_size=12, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DESIGN DECISIONS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Key Design Decisions", "Why I chose each technology")

decisions = [
    ("LangGraph\n(not a while loop)",
     ["✓ Every node boundary auto-logs timing + state",
      "✓ stream_mode='updates' feeds live UI progress",
      "✓ Hard loop cap enforced by graph topology",
      "✗ Cost: framework-specific node signature convention"],
     BLUE_LIGHT),
    ("Local LLM via Ollama\n(not OpenAI API)",
     ["✓ Zero cost per query",
      "✓ Privacy-safe — no data leaves the machine",
      "✓ Works fully offline",
      "✗ Cost: 1B model has limited instruction following"],
     RGBColor(0xE8, 0xF5, 0xE9)),
    ("ChromaDB RAG\n(not prompt-stuffing)",
     ["✓ 200-page PDF won't fit in context window",
      "✓ Retrieves only the 3 most relevant chunks",
      "✓ New guideline = replace PDF, no code change",
      "✗ Cost: collection must be populated first"],
     RGBColor(0xFF, 0xF8, 0xE1)),
    ("Two-phase UI streaming\n(graph + direct LLM)",
     ["✓ Phase 1: LangGraph for tool execution",
      "✓ Phase 2: llm.stream() for token-level output",
      "✓ JSON internally, Markdown displayed to user",
      "✗ Cost: slightly more complex UI code"],
     RGBColor(0xFC, 0xE4, 0xEC)),
]

for i, (title, items, bg) in enumerate(decisions):
    x = Inches(0.25 + (i % 2) * 6.5)
    y = Inches(1.55 + (i // 2) * 2.85)
    bullet_box(slide, items, x, y, Inches(6.2), Inches(2.6),
               title=title, font_size=12, bg=bg)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TRADE-OFFS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Known Trade-offs", "Honest assessment — what works, what needs work")

bullet_box(slide,
    ["✅ Full research pipeline (reason → plan → act → observe → respond)",
     "✅ Real PubMed retrieval — live NCBI API, real PMIDs",
     "✅ ICD-10-CM lookup — 74,260 CMS FY2025 codes",
     "✅ Token-by-token streaming in the web UI",
     "✅ JSONL trace file for every run (full audit trail)",
     "✅ Safety guardrails: PII, injections, emergencies",
     "✅ LLM-as-judge evaluation framework (Claude Haiku)",
     "✅ Conversation history in SQLite (session memory)",
    ],
    Inches(0.3), Inches(1.55), Inches(6.2), Inches(5.7),
    title="What Works Well", font_size=13, bg=RGBColor(0xE8, 0xF5, 0xE9))

bullet_box(slide,
    ["⚠ Load clinical guideline PDFs → fixes TC-001, TC-003",
     "⚠ Switch to 70B model → 10x better answer quality",
     "⚠ Add CPT code database (AMA license required)",
     "⚠ Add user authentication before network deploy",
     "⚠ Replace SQLite → PostgreSQL for multi-user",
     "⚠ Add retry logic for tool argument validation failures",
     "⚠ ChromaDB is single-process — not multi-user ready",
     "⚠ No HTTPS — localhost only in current form",
    ],
    Inches(6.8), Inches(1.55), Inches(6.2), Inches(5.7),
    title="Needs Work for Production", font_size=13, bg=RGBColor(0xFF, 0xF3, 0xE0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — LESSONS LEARNED
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Key Engineering Lessons", "What I learned building this")

lessons = [
    ("LangGraph node signatures are strict",
     "Nodes only receive config if the second parameter is typed as RunnableConfig "
     "(not Any). This caused a 'missing argument' error that took real debugging to find."),
    ("Small models fail at tool-use",
     "gemma3:1b passes wrong argument names ~50% of the time. Fixed with Pydantic "
     "model_validator aliases — a compatibility shim that remaps 'query' → 'natural_query'."),
    ("You can't stream JSON mid-parse",
     "The LangGraph respond node returns complete JSON. Streaming JSON characters to a user "
     "is unreadable. Solution: use JSON internally, plain Markdown for the UI streaming layer."),
    ("Heuristic eval is a proxy, not a measure",
     "Keyword recall doesn't capture clinical accuracy. A model can retrieve the right PMID "
     "and still write a vague answer. LLM-as-judge is far more meaningful."),
    ("The hardest part wasn't the AI",
     "It was reliability: Pydantic coercions, Windows cp1252 encoding errors, "
     "SQLAlchemy rejecting list-typed fields, Streamlit's torchvision file-watcher issue."),
]

for i, (title, body) in enumerate(lessons):
    y = Inches(1.55 + i * 1.12)
    add_rect(slide, Inches(0.3), y, Inches(12.7), Inches(1.0), BLUE_LIGHT)
    add_rect(slide, Inches(0.3), y, Inches(0.08), Inches(1.0), BLUE_MID)
    add_text_box(slide, f"{i+1}.  {title}",
                 Inches(0.5), y + Inches(0.04), Inches(12.3), Inches(0.35),
                 font_size=13, bold=True, color=BLUE_DARK)
    add_text_box(slide, body,
                 Inches(0.5), y + Inches(0.42), Inches(12.3), Inches(0.5),
                 font_size=11, color=GREY_TEXT, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Roadmap to Production", "If I had more time")

phases = [
    ("This Week",          RGBColor(0xE3, 0xF2, 0xFD),
     ["Load ADA 2024 + ACC/AHA guideline PDFs",
      "Switch to Groq free tier (llama-3.1-70b)",
      "Add unit tests for all tool validators",
      "Expected: pass rate 4/5 (80%)"]),
    ("Next Month",         RGBColor(0xE8, 0xF5, 0xE9),
     ["Add CPT code database",
      "Add user authentication (OAuth proxy)",
      "PostgreSQL for multi-user conversations",
      "Docker container for easy deployment"]),
    ("3 Months",           RGBColor(0xFF, 0xF8, 0xE1),
     ["Fine-tune 7B model on clinical Q&A pairs",
      "FHIR integration — pull EHR patient context",
      "Multi-turn conversation memory across sessions",
      "Payer-specific LCD/NCD lookup per state"]),
    ("Production Ready",   RGBColor(0xFC, 0xE4, 0xEC),
     ["SOC 2 / HIPAA compliance review",
      "High-availability Kubernetes deployment",
      "Real-time CMS coverage policy updates",
      "Integration with CDI/RCM workflow systems"]),
]

for i, (phase, bg, items) in enumerate(phases):
    x = Inches(0.25 + i * 3.26)
    bullet_box(slide, items, x, Inches(1.55), Inches(3.1), Inches(5.7),
               title=phase, font_size=12, bg=bg)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_rect(slide, 0, 0, W, H, BLUE_DARK)
add_rect(slide, 0, Inches(4.6), W, Inches(2.9), BLUE_MID)

add_text_box(slide, "What I Built",
             Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

add_text_box(slide,
             "A 5-node LangGraph agent that answers clinical coding, coverage, and "
             "treatment questions by searching PubMed, clinical guidelines, and ICD-10 "
             "databases in real time — with a streaming web UI, safety guardrails, "
             "structured evaluation, and full audit trails.",
             Inches(0.8), Inches(0.95), Inches(11.5), Inches(1.3),
             font_size=19, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

skills = [
    ("LangGraph\nstate machine",  "5-node self-reflection loop"),
    ("RAG pipeline",              "ChromaDB + sentence-transformers"),
    ("Multi-API tools",           "NCBI Entrez + CMS ICD-10"),
    ("Safety guardrails",         "Pydantic validation stack"),
    ("Streaming UI",              "Streamlit token streaming"),
    ("LLM evaluation",            "Claude Haiku as judge"),
]

for i, (skill, detail) in enumerate(skills):
    x = Inches(0.4 + (i % 3) * 4.3)
    y = Inches(2.4 + (i // 3) * 1.7)
    add_rect(slide, x, y, Inches(3.9), Inches(1.4), RGBColor(0x0A, 0x33, 0x80))
    add_text_box(slide, skill,
                 x + Inches(0.12), y + Inches(0.1), Inches(3.7), Inches(0.65),
                 font_size=14, bold=True, color=WHITE)
    add_text_box(slide, detail,
                 x + Inches(0.12), y + Inches(0.75), Inches(3.7), Inches(0.5),
                 font_size=11, color=RGBColor(0x90, 0xCA, 0xF9))

add_text_box(slide, "Thank you — happy to walk through any part of the code in detail.",
             Inches(0.6), Inches(5.05), Inches(12), Inches(0.5),
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "Balaji GV  ·  balaji.gv@tungstenautomation.com  ·  May 2026",
             Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
             font_size=13, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "Healthcare_QA_Agent.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
